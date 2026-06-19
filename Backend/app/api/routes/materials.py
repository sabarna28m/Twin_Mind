import base64
import json
import logging
import os
import uuid
from io import BytesIO
from pathlib import Path
from typing import List, Optional

import requests as http_requests
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, status
from fastapi.responses import FileResponse
from pydantic import BaseModel
from sqlalchemy.orm import Session as DBSession

from app.core.config import settings
from app.core.database import get_db
from app.models.material import Material
from app.models.user import User
from app.api.routes.auth import get_current_user
from app.api.schemas.material import MaterialResponse

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/materials", tags=["materials"])

UPLOAD_DIR = Path(__file__).resolve().parents[3] / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

MAX_FILE_SIZE = 20 * 1024 * 1024  # 20 MB
ALLOWED_TYPES = {
    "application/pdf",
    "image/png", "image/jpeg", "image/gif", "image/webp",
    "text/plain",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/msword",
}

GROQ_MODEL = "llama-3.3-70b-versatile"


# ── Text extraction ────────────────────────────────────────────────────────────

def _extract_text(file_bytes: bytes, mime_type: str, max_chars: int = 12000) -> str:
    try:
        if mime_type.startswith("text/"):
            return file_bytes.decode("utf-8", errors="replace")[:max_chars]

        if mime_type == "application/pdf":
            from pypdf import PdfReader
            reader = PdfReader(BytesIO(file_bytes))
            pages = [page.extract_text() or "" for page in reader.pages[:30]]
            return "\n".join(pages)[:max_chars]

        if "wordprocessingml" in mime_type or "msword" in mime_type:
            from docx import Document
            doc = Document(BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:max_chars]

        if "presentationml" in mime_type:
            try:
                from pptx import Presentation
                prs = Presentation(BytesIO(file_bytes))
                texts = [
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                return "\n".join(texts)[:max_chars]
            except ImportError:
                logger.warning("python-pptx not installed — PPTX text extraction unavailable")
                return ""

    except Exception as exc:
        logger.warning("Text extraction failed for %s: %s", mime_type, exc)
    return ""


# ── Gemini image description ───────────────────────────────────────────────────

def _describe_with_gemini(image_bytes: bytes, mime_type: str, prompt: str) -> str:
    if not settings.gemini_api_key:
        return ""
    url = (
        "https://generativelanguage.googleapis.com/v1beta/models"
        f"/gemini-2.0-flash:generateContent?key={settings.gemini_api_key}"
    )
    body = {
        "contents": [{
            "parts": [
                {"inline_data": {"mime_type": mime_type, "data": base64.b64encode(image_bytes).decode()}},
                {"text": prompt},
            ]
        }]
    }
    try:
        resp = http_requests.post(url, json=body, timeout=45)
        if resp.ok:
            return resp.json()["candidates"][0]["content"]["parts"][0]["text"]
    except Exception as exc:
        logger.warning("Gemini error: %s", exc)
    return ""


# ── Groq JSON call ─────────────────────────────────────────────────────────────

def _call_groq(prompt: str, system: str, max_tokens: int = 3500) -> dict:
    if not settings.groq_api_key:
        raise HTTPException(status_code=503, detail="AI service not configured (GROQ_API_KEY missing)")
    from groq import Groq
    client = Groq(api_key=settings.groq_api_key)
    try:
        resp = client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {"role": "system", "content": system},
                {"role": "user",   "content": prompt},
            ],
            max_tokens=max_tokens,
            temperature=0.3,
            response_format={"type": "json_object"},
        )
        raw = resp.choices[0].message.content or "{}"
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise HTTPException(status_code=500, detail=f"AI returned invalid JSON: {exc}")
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"AI call failed: {exc}")


# ── Helper: load + extract material content ────────────────────────────────────

def _load_material_text(record: Material, image_prompt: str = "Describe all educational content in detail.") -> str:
    path = UPLOAD_DIR / str(record.user_id) / record.stored_name
    if not path.exists():
        raise HTTPException(status_code=404, detail="File missing from storage")
    file_bytes = path.read_bytes()

    if record.mime_type.startswith("image/"):
        text = _describe_with_gemini(file_bytes, record.mime_type, image_prompt)
        if not text:
            raise HTTPException(
                status_code=503,
                detail="Image analysis requires GEMINI_API_KEY to be configured."
            )
        return text

    text = _extract_text(file_bytes, record.mime_type)
    if not text:
        raise HTTPException(
            status_code=422,
            detail=f"Could not extract text from '{record.original_name}'. "
                   "Ensure the file contains readable text."
        )
    return text


# ── CRUD endpoints ─────────────────────────────────────────────────────────────

@router.get("", response_model=List[MaterialResponse])
def list_materials(
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    return (
        db.query(Material)
        .filter(Material.user_id == current_user.id)
        .order_by(Material.created_at.desc())
        .all()
    )


@router.post("", response_model=MaterialResponse, status_code=status.HTTP_201_CREATED)
async def upload_material(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE, detail="File exceeds 20 MB limit")

    mime = file.content_type or "application/octet-stream"
    if mime not in ALLOWED_TYPES:
        raise HTTPException(status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE, detail=f"File type '{mime}' not allowed")

    ext = Path(file.filename or "file").suffix
    stored_name = f"{uuid.uuid4().hex}{ext}"
    user_dir = UPLOAD_DIR / str(current_user.id)
    user_dir.mkdir(exist_ok=True)
    (user_dir / stored_name).write_bytes(content)

    record = Material(
        user_id=current_user.id,
        original_name=file.filename or stored_name,
        stored_name=stored_name,
        mime_type=mime,
        file_size=len(content),
    )
    db.add(record)
    db.commit()
    db.refresh(record)
    return record


@router.get("/{material_id}/download")
def download_material(
    material_id: int,
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    record = db.query(Material).filter(Material.id == material_id, Material.user_id == current_user.id).first()
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Material not found")
    path = UPLOAD_DIR / str(current_user.id) / record.stored_name
    if not path.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File missing from storage")
    return FileResponse(path=str(path), filename=record.original_name, media_type=record.mime_type)


@router.delete("/{material_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_material(
    material_id: int,
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    record = db.query(Material).filter(Material.id == material_id, Material.user_id == current_user.id).first()
    if not record:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Material not found")
    path = UPLOAD_DIR / str(current_user.id) / record.stored_name
    if path.exists():
        path.unlink()
    db.delete(record)
    db.commit()


# ── AI Analysis endpoints ──────────────────────────────────────────────────────

@router.post("/{material_id}/analyze")
async def analyze_material(
    material_id: int,
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    """Comprehensive AI analysis: summary, flashcards, questions, definitions, mind map."""
    record = db.query(Material).filter(Material.id == material_id, Material.user_id == current_user.id).first()
    if not record:
        raise HTTPException(status_code=404, detail="Material not found")

    text = _load_material_text(
        record,
        "Describe all text, diagrams, concepts, and educational content in this image in full detail."
    )

    system = (
        "You are an expert educational content analyzer. "
        "Respond only with a single valid JSON object. No markdown, no code blocks."
    )
    prompt = f"""Analyze this educational document titled "{record.original_name}" and return a JSON object with EXACTLY these keys:

"quick_summary": [6-8 concise bullet point strings covering the main points]
"detailed_summary": "2-3 paragraph narrative summary of the entire document"
"chapter_breakdown": [array of {{"title": "...", "summary": "...", "key_points": ["..."]}} per chapter/section]
"important_concepts": [up to 12 concept name strings]
"key_definitions": [up to 15 {{"term": "...", "definition": "..."}} objects]
"exam_notes": "string of focused notes on what is most likely to be tested"
"flashcards": [up to 20 {{"front": "question or term", "back": "answer or definition"}} objects]
"important_questions": [up to 15 likely exam question strings]
"revision_notes": "concise 200-word revision summary"
"topics": [up to 10 main topic strings]
"keywords": [up to 25 key term strings]
"mind_map": {{"central": "main topic string", "branches": [{{"label": "branch topic", "children": ["sub-topic string"]}}]}}

Document content:
{text[:10000]}"""

    result = _call_groq(prompt, system, max_tokens=4500)
    result["material_id"] = material_id
    result["material_name"] = record.original_name
    return result


@router.post("/{material_id}/analyze-paper")
async def analyze_paper(
    material_id: int,
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    """Analyze an exam/question paper: pattern, marks distribution, topic weightage."""
    record = db.query(Material).filter(Material.id == material_id, Material.user_id == current_user.id).first()
    if not record:
        raise HTTPException(status_code=404, detail="Material not found")

    text = _load_material_text(
        record,
        "Extract every question, instruction, marks allocation, section heading, and time limit from this question paper image."
    )

    system = "You are an expert exam paper analyzer. Respond only with a single valid JSON object."
    prompt = f"""Analyze this exam/question paper titled "{record.original_name}" and return a JSON object with EXACTLY these keys:

"exam_pattern": "string describing the overall pattern and structure"
"total_marks": number (0 if unknown)
"total_time": "string like '3 hours' or 'unknown'"
"total_questions": number (0 if unknown)
"marks_distribution": [{{"section": "A", "total_marks": 20, "questions_count": 4, "marks_per_question": 5}}]
"section_structure": [{{"section": "A", "title": "...", "question_type": "MCQ|Short Answer|Essay|etc.", "instructions": "..."}}]
"question_types": ["MCQ", "Short Answer", "Essay"] (all types present)
"difficulty_level": "Easy" or "Medium" or "Hard" or "Mixed"
"difficulty_distribution": {{"easy": 30, "medium": 50, "hard": 20}} (percentages summing to 100)
"topic_weightage": [{{"topic": "...", "percentage": 25, "questions_count": 5}}]
"frequently_repeated_topics": ["topic likely to appear again"]
"question_format": "description of how questions are phrased and formatted"
"assessment_style": "description of the overall assessment approach"

Paper content:
{text[:10000]}"""

    result = _call_groq(prompt, system, max_tokens=3000)
    result["material_id"] = material_id
    result["material_name"] = record.original_name
    return result


class GeneratePaperRequest(BaseModel):
    analysis: dict
    difficulty: str = "same"  # same | easier | harder | mixed


@router.post("/{material_id}/generate-paper")
async def generate_paper(
    material_id: int,
    body: GeneratePaperRequest,
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    """Generate a new question paper following the same pattern but with entirely new questions."""
    record = db.query(Material).filter(Material.id == material_id, Material.user_id == current_user.id).first()
    if not record:
        raise HTTPException(status_code=404, detail="Material not found")

    difficulty_note = {
        "same":   "Match the EXACT same difficulty level as the original.",
        "easier": "Make it EASIER than the original — simpler questions, more straightforward.",
        "harder": "Make it HARDER than the original — more complex questions, deeper reasoning required.",
        "mixed":  "Use MIXED difficulty: approximately 30% easy, 50% medium, 20% hard.",
    }.get(body.difficulty, "Match the same difficulty as the original.")

    analysis_json = json.dumps(body.analysis)[:3500]
    system = "You are an expert exam paper generator. Respond only with a single valid JSON object."
    prompt = f"""Based on this exam paper pattern analysis:
{analysis_json}

Generate a BRAND NEW question paper following the SAME structure, sections, marks distribution, and question types.
Do NOT copy or paraphrase any question from the original paper. All questions must be entirely new.

Difficulty: {difficulty_note}

Return a JSON object with EXACTLY these keys:
"title": "Generated Question Paper — [subject/topic]"
"difficulty": "{body.difficulty}"
"total_marks": number
"total_time": "string"
"instructions": "general exam instructions string"
"sections": [{{
  "section": "A",
  "title": "section title",
  "instructions": "section-specific instructions",
  "questions": [{{
    "number": 1,
    "question": "question text",
    "type": "MCQ|Short Answer|Essay|Fill in the Blank|True/False",
    "options": ["A) ...", "B) ...", "C) ...", "D) ..."] (ONLY for MCQ, omit for other types),
    "marks": number
  }}]
}}]"""

    result = _call_groq(prompt, system, max_tokens=4500)
    result["material_id"] = material_id
    result["difficulty"] = body.difficulty
    return result


class GenerateQuizRequest(BaseModel):
    subject: Optional[str] = ""
    difficulty: str = "medium"
    count: int = 10
    question_types: List[str] = ["MCQ"]


@router.post("/{material_id}/generate-quiz")
async def generate_quiz_from_material(
    material_id: int,
    body: GenerateQuizRequest,
    current_user: User = Depends(get_current_user),
    db: DBSession = Depends(get_db),
):
    """Generate an interactive quiz directly from an uploaded study material."""
    record = db.query(Material).filter(Material.id == material_id, Material.user_id == current_user.id).first()
    if not record:
        raise HTTPException(status_code=404, detail="Material not found")

    text = _load_material_text(
        record,
        "Describe all educational content, concepts, facts, and text in this study material image."
    )

    subject_str  = f" on {body.subject}" if body.subject else ""
    types_str    = ", ".join(body.question_types)
    count        = max(1, min(body.count, 30))

    system = "You are an expert quiz generator for students. Respond only with a single valid JSON object."
    prompt = f"""Create a quiz{subject_str} from this study material titled "{record.original_name}".

Requirements:
- Exactly {count} questions
- Difficulty: {body.difficulty}
- Question types to include: {types_str}

Return a JSON object with a single key "questions" containing an array of question objects.
Each question object MUST have ALL of these keys:
- "id": integer (1-based)
- "type": one of "MCQ" | "true_false" | "fill_blank" | "short_answer" | "long_answer"
- "question": the question text string
- "options": array of 4 strings like ["A) ...", "B) ...", "C) ...", "D) ..."] — ONLY for MCQ type, OMIT for all other types
- "correct_answer": for MCQ use the letter "A", "B", "C", or "D"; for true_false use "True" or "False"; for fill_blank/short_answer/long_answer use the full answer string
- "explanation": brief explanation string (1-2 sentences) explaining why the answer is correct
- "difficulty": "easy" | "medium" | "hard"
- "topic": string naming the specific sub-topic this question tests
- "marks": 1 for easy/medium, 2 for hard

Material content:
{text[:9000]}"""

    result = _call_groq(prompt, system, max_tokens=4500)
    result["material_id"] = material_id
    result["material_name"] = record.original_name
    result["difficulty"] = body.difficulty
    result["question_types"] = body.question_types
    return result
